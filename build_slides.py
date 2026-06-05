"""
build_slides.py — Generate the actual PowerPoint file for the team meeting.

Outputs: team_meeting.pptx (16:9 widescreen, ~9 slides + title + Q&A backup)

Run:
    python build_slides.py
    python build_slides.py --output team_meeting.pptx
"""

import argparse
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------- Colors ----------
NAVY    = RGBColor(0x1F, 0x3A, 0x5F)   # title / accent
DARK    = RGBColor(0x2C, 0x3E, 0x50)   # body text
GRAY    = RGBColor(0x7F, 0x8C, 0x8D)   # subtitles
ORANGE  = RGBColor(0xF3, 0x9C, 0x12)   # highlight (methodology)
PURPLE  = RGBColor(0x9B, 0x59, 0xB6)   # HPS
GREEN   = RGBColor(0x27, 0xAE, 0x60)   # C4 / wins
RED     = RGBColor(0xE7, 0x4C, 0x3C)   # critical / attack
BLUE    = RGBColor(0x34, 0x98, 0xDB)   # benign / info
LIGHT   = RGBColor(0xEC, 0xF0, 0xF1)   # background fill
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)


# ---------- Geometry helpers (16:9 standard) ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.5)
MARGIN_Y = Inches(0.5)


def add_title(slide, text, color=NAVY, size=32):
    box = slide.shapes.add_textbox(MARGIN_X, MARGIN_Y, SLIDE_W - 2 * MARGIN_X, Inches(0.7))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return box


def add_subtitle(slide, text, top, color=GRAY, size=16):
    box = slide.shapes.add_textbox(MARGIN_X, top, SLIDE_W - 2 * MARGIN_X, Inches(0.4))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.italic = True
    r.font.color.rgb = color
    return box


def add_textbox(slide, text, left, top, width, height,
                size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    if isinstance(text, str):
        lines = text.split("\n")
    else:
        lines = text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def add_bullets(slide, items, left, top, width, height, size=13, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        r = p.add_run()
        r.text = "• " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def add_box(slide, left, top, width, height, fill_color, line_color=None):
    """Add a colored rectangular shape (used for callouts)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_callout(slide, text, left, top, width, height,
                fill=LIGHT, border=ORANGE, text_color=DARK, size=14, bold=False):
    add_box(slide, left, top, width, height, fill, border)
    tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.05),
                                  width - Inches(0.2), height - Inches(0.1))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color


def add_image(slide, path, left, top, width=None, height=None):
    """Add image if file exists, else a placeholder box with the missing path."""
    if path and os.path.exists(path):
        if width and height:
            return slide.shapes.add_picture(path, left, top, width=width, height=height)
        elif width:
            return slide.shapes.add_picture(path, left, top, width=width)
        else:
            return slide.shapes.add_picture(path, left, top)
    # Placeholder box
    w = width or Inches(5)
    h = height or Inches(3)
    add_box(slide, left, top, w, h, LIGHT, GRAY)
    add_textbox(slide, [f"[ figure missing ]", path or "(no path)"],
                left, top, w, h, size=11, color=GRAY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_table(slide, headers, rows, left, top, width, height,
              header_color=NAVY, header_text=WHITE, body_size=12,
              bold_first_col=False, highlight_rows=None):
    """Add a styled table. highlight_rows is a list of row indices to highlight in red."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2);  tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.text = ""
        r = p.add_run()
        r.text = h
        r.font.size = Pt(body_size)
        r.font.bold = True
        r.font.color.rgb = header_text

    # Body rows
    highlight_rows = highlight_rows or []
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            if i in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(4); tf.margin_right = Pt(4)
            tf.margin_top = Pt(2);  tf.margin_bottom = Pt(2)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.text = ""
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(body_size)
            r.font.color.rgb = DARK
            if bold_first_col and j == 0:
                r.font.bold = True
            if i in highlight_rows:
                r.font.bold = True
                r.font.color.rgb = RED
    return tbl


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


# ============================================================
#  SLIDE 0: TITLE
# ============================================================
def slide_title(prs):
    s = blank_slide(prs)
    # Background accent
    add_box(s, Inches(0), Inches(0), SLIDE_W, Inches(1.5), NAVY)

    add_textbox(s, "Hyperbolic Geometric Priors\nfor LLM Jailbreak Detection",
                Inches(0.5), Inches(2.2), SLIDE_W - Inches(1), Inches(1.5),
                size=44, color=NAVY, bold=True)

    add_textbox(s, "A Controlled Comparison Study",
                Inches(0.5), Inches(3.6), SLIDE_W - Inches(1), Inches(0.5),
                size=22, color=GRAY)

    add_textbox(s, ["Methodology critique • HPS vs C4 vs MTP",
                    "Cross-model evaluation: Llama-3-8B vs Vicuna-13B"],
                Inches(0.5), Inches(4.5), SLIDE_W - Inches(1), Inches(1.0),
                size=16, color=DARK)

    add_textbox(s, "Team meeting — Findings & Open Questions",
                Inches(0.5), Inches(6.5), SLIDE_W - Inches(1), Inches(0.5),
                size=14, color=GRAY)


# ============================================================
#  SLIDE 1: THE QUESTION
# ============================================================
def slide_1(prs, fig_dir):
    s = blank_slide(prs)
    add_title(s, "1. Do Hyperbolic Priors Help Jailbreak Detection?")
    add_subtitle(s, "A controlled comparison study", Inches(1.05))

    # Left: motivation bullets
    add_textbox(s, "Background", Inches(0.5), Inches(1.6), Inches(6), Inches(0.4),
                size=16, color=NAVY, bold=True)
    add_bullets(s, [
        "HypLoRA (NeurIPS 2025): LLM token embeddings show empirical δ-hyperbolicity (tree-like structure)",
        "HELM (NeurIPS 2025): Power-law radial structure / negative Ricci curvature",
        "Theoretical prediction: hyperbolic projection should provide useful inductive bias for distinguishing harmful (specific) from benign (general)",
    ], Inches(0.5), Inches(2.0), Inches(6.3), Inches(2.5), size=13)

    # Research question callout
    add_callout(s,
                '"If LLM activations have hierarchical structure, does '
                "hyperbolic projection provide a measurable detection advantage "
                'over flat (Euclidean) baselines?"',
                Inches(0.5), Inches(4.7), Inches(6.3), Inches(1.4),
                fill=RGBColor(0xFE, 0xF9, 0xE7), border=ORANGE, size=14, bold=True)

    # Three deliverables at the bottom
    add_textbox(s, "What we built:", Inches(0.5), Inches(6.2), Inches(6.3), Inches(0.3),
                size=13, color=NAVY, bold=True)
    add_bullets(s, [
        "HPS (Hyperbolic Projection Sentinel): Lorentz contrastive framework",
        "Parameter-matched Euclidean ablation + linear probe baseline",
        "Two LLMs: Llama-3-8B-Instruct (RLHF) vs Vicuna-13B-v1.5 (SFT-only)",
    ], Inches(0.5), Inches(6.5), Inches(6.3), Inches(1.0), size=11)

    # Right: concept figure
    add_image(s, os.path.join(fig_dir, "fig_lorentz_concept.png"),
              Inches(7.0), Inches(1.6), width=Inches(6.0))
    add_textbox(s, "Hypothesis: attacks at higher radial position",
                Inches(7.0), Inches(6.7), Inches(6.0), Inches(0.3),
                size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
#  SLIDE 2: METHODS
# ============================================================
def slide_2(prs, fig_dir):
    s = blank_slide(prs)
    add_title(s, "2. Methods Under Comparison")
    add_subtitle(s, "Four-method controlled comparison with parameter-matched ablation",
                 Inches(1.05))

    headers = ["Method", "Architecture", "Parameters", "Innovation"]
    rows = [
        ["HPS (ours)", "Lorentz proj + 12 trajectory features + LR", "262K",
         "Hyperbolic geometric prior"],
        ["HPS-Euclidean (control)", "Same arch, FLAT geometry", "262K (matched)",
         "Tests if geometry matters"],
        ["C4 (baseline)", "Mean-pool 6 layers' last-token + LR", "4,097",
         "Adapted from Anthropic Cheap Monitors"],
        ["MTP (baseline)", "Mean-pool tokens at 1 layer + LR", "4,097",
         "Faithful Anthropic reproduction"],
    ]
    add_table(s, headers, rows,
              Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.0),
              body_size=12, bold_first_col=True)

    # HPS architecture diagram
    add_callout(s,
                "HPS pipeline: Activations [N=6 layers] → W·activations (project to 64-d) → "
                "Lorentz hyperboloid (κ=0.1, frozen) → 12 trajectory features "
                "(5 radial + 4 curvature + 3 displacement) → LR.\n"
                "Trained with per-layer-temperature contrastive loss, 50 epochs.",
                Inches(0.5), Inches(4.0), Inches(12.3), Inches(1.2),
                fill=LIGHT, border=NAVY, size=12)

    # Why the parameter-matched control matters
    add_callout(s,
                "Why HPS-Euclidean control is critical:\n"
                "Without parameter-matched Euclidean, any HPS-vs-C4 win could be 'more parameters helped.' "
                "HPS-Euclidean has IDENTICAL parameter count and architecture, only flat geometry. "
                "This isolates the geometric prior.",
                Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.6),
                fill=RGBColor(0xFE, 0xF9, 0xE7), border=ORANGE, size=13, bold=True)


# ============================================================
#  SLIDE 3: LENGTH CONFOUND
# ============================================================
def slide_3(prs, fig_dir):
    s = blank_slide(prs)
    add_title(s, "3. Methodology Surprise #1 — Length Confound", color=ORANGE)
    add_subtitle(s,
                 "A trivial classifier using only prompt length detects 96.9% of attacks",
                 Inches(1.05))

    # Left: the headline numbers
    add_textbox(s, "Length-only classifier (just len(prompt), no model)",
                Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.4),
                size=14, color=NAVY, bold=True)

    headers = ["Metric", "Length-only"]
    rows = [
        ["AUROC", "0.992"],
        ["TPR @ 5% FPR", "0.969"],
        ["TPR @ 1% FPR", "0.875"],
    ]
    add_table(s, headers, rows,
              Inches(0.5), Inches(2.0), Inches(5.5), Inches(1.4),
              body_size=14, bold_first_col=True)

    add_callout(s,
                "A classifier that NEVER LOOKS at the prompt CONTENT — only its length — "
                "catches 96.9% of attacks at 5% FPR.",
                Inches(0.5), Inches(3.6), Inches(5.5), Inches(1.0),
                fill=RGBColor(0xFE, 0xF9, 0xE7), border=ORANGE, size=12, bold=True)

    # Length distribution
    add_textbox(s, "Length distribution",
                Inches(0.5), Inches(4.8), Inches(5.5), Inches(0.3),
                size=13, color=NAVY, bold=True)
    add_bullets(s, [
        "Alpaca harmless: very short (mean ~12-60 chars)",
        "Attacks: weighted mean 218 chars, extreme by-attack variation",
        "Per-attack lengths range 35 (gcg) → 2,193 (puzzler)",
    ], Inches(0.5), Inches(5.1), Inches(5.5), Inches(1.5), size=11)

    # Right: per-attack length AUROC
    add_textbox(s, "Per-attack length-only AUROC",
                Inches(6.5), Inches(1.6), Inches(6.3), Inches(0.4),
                size=14, color=NAVY, bold=True)

    headers2 = ["Attack", "Mean length", "Length-only AUROC"]
    rows2 = [
        ["puzzler", "2,193 chars", "1.000"],
        ["saa", "473 chars", "1.000"],
        ["drattack", "437 chars", "1.000"],
        ["base64", "124 chars", "1.000"],
        ["ijp", "460 chars", "0.999"],
        ["autodan", "77 chars", "0.999"],
        ["pair", "72 chars", "0.997"],
        ["zulu", "50 chars", "0.976"],
        ["gcg", "35 chars", "0.971"],
    ]
    add_table(s, headers2, rows2,
              Inches(6.5), Inches(2.0), Inches(6.3), Inches(3.6),
              body_size=11, bold_first_col=True)

    # Bottom: the fix
    add_callout(s,
                "Our fix: Diverse benign data (WildChat + OR-Bench + MMLU + GSM8K + WikiText long-form). "
                "After fix → length-only AUROC: 0.992 → 0.318. Permutation: 0.498 (real signal exists).",
                Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.9),
                fill=RGBColor(0xE9, 0xF7, 0xEF), border=GREEN, size=12, bold=True)


# ============================================================
#  SLIDE 4: max_length CONFOUND
# ============================================================
def slide_4(prs, fig_dir):
    s = blank_slide(prs)
    add_title(s, "4. Methodology Surprise #2 — max_length Confound", color=ORANGE)
    add_subtitle(s,
                 "Inconsistent token truncation → norm-only AUROC = 1.000",
                 Inches(1.05))

    # Left: the bug code block
    add_textbox(s, "The bug:", Inches(0.5), Inches(1.6), Inches(6), Inches(0.3),
                size=14, color=NAVY, bold=True)
    add_callout(s,
                "Original cache:\n"
                "  benign max_length = 512    attacks max_length = 512   ✓ consistent\n"
                "  → Norm-only AUROC = 0.917\n\n"
                "Diverse benign cache (after length confound fix):\n"
                "  benign max_length = 2048   attacks REUSED at 512   ✗ MISMATCH\n"
                "  → Norm-only AUROC = 1.000  ← perfect detection from norm alone",
                Inches(0.5), Inches(2.0), Inches(6.3), Inches(2.3),
                fill=LIGHT, border=DARK, size=11, bold=False)

    # Layer 31 norm comparison
    add_textbox(s, "Layer 31 (deepest) activation norms:",
                Inches(0.5), Inches(4.5), Inches(6.3), Inches(0.3),
                size=13, color=NAVY, bold=True)

    headers = ["", "Old cache (consistent)", "Diverse cache (mismatch)"]
    rows = [
        ["Benign mean norm", "155.6", "35.5  ⚠"],
        ["Attack mean norm", "153.0", "153.0"],
        ["Ratio", "1.0×", "4.3×"],
    ]
    add_table(s, headers, rows,
              Inches(0.5), Inches(4.9), Inches(6.3), Inches(1.5),
              body_size=12, bold_first_col=True, highlight_rows=[2])

    # Bottom callout
    add_callout(s,
                "The fix: Re-extract attacks at max_length=2048 to match diverse benign.\n"
                "  Norm-only AUROC: 1.000 → 0.761 ✓ (confound largely resolved)\n"
                "  Permutation test AUROC = 0.498 → real signal exists",
                Inches(0.5), Inches(6.5), Inches(6.3), Inches(0.85),
                fill=RGBColor(0xE9, 0xF7, 0xEF), border=GREEN, size=11, bold=True)

    # Right: figure (norm-controlled eval not pulled, leave space)
    add_box(s, Inches(7.2), Inches(1.6), Inches(5.8), Inches(5.4), LIGHT, GRAY)
    add_textbox(s,
                ["Norm-controlled evaluation",
                 "",
                 "Across L2-normalize, standardize, and",
                 "L2+std controls:",
                 "",
                 "  norm-only AUROC:  0.76 → 0.55",
                 "  C4 TPR @ 5% FPR:  0.995 → 0.996",
                 "",
                 "→ C4 is robust to all norm controls.",
                 "  Real semantic signal exists beyond norm."],
                Inches(7.4), Inches(1.8), Inches(5.4), Inches(5.0),
                size=14, color=DARK, bold=False)
    add_textbox(s, "(see paper appendix for full norm-control table)",
                Inches(7.2), Inches(7.0), Inches(5.8), Inches(0.3),
                size=9, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
#  SLIDE 5: GEOMETRIC HYPOTHESIS CONFIRMED
# ============================================================
def slide_5(prs, fig_dir):
    s = blank_slide(prs)
    add_title(s, "5. Finding 1 — Geometric Hypothesis Confirmed", color=GREEN)
    add_subtitle(s,
                 "After methodology fixes: 0/13 inversions across all configurations",
                 Inches(1.05))

    # Left: the test
    add_textbox(s, "13 configurations tested",
                Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.3),
                size=14, color=NAVY, bold=True)
    add_bullets(s, [
        "5 random seeds (κ=0.1, 50 epochs)",
        "4 epoch checkpoints (5, 10, 25, 50)",
        "4 curvature κ values (0.1, 0.5, 1.0, 2.0)",
    ], Inches(0.5), Inches(2.0), Inches(5.5), Inches(1.2), size=12)

    add_callout(s,
                "Result: 0/13 inversions.\n"
                "All configurations show benign median < attack median.\n"
                "Hyperbolic prior is VALIDATED.",
                Inches(0.5), Inches(3.3), Inches(5.5), Inches(1.4),
                fill=RGBColor(0xE9, 0xF7, 0xEF), border=GREEN,
                size=14, bold=True)

    # Mechanism
    add_textbox(s, "Mechanism:", Inches(0.5), Inches(4.9), Inches(5.5), Inches(0.3),
                size=13, color=NAVY, bold=True)
    add_textbox(s,
                ["Hyperbolic prior predicts: harmful (specific) → high radial.",
                 "Before length-confound fix: appeared inverted (length-shortcut).",
                 "After fix: attacks ARE at higher radial, as theory predicts."],
                Inches(0.5), Inches(5.2), Inches(5.5), Inches(1.5), size=11)

    add_callout(s,
                "The geometric prior is theoretically valid. "
                "Hyperbolic projection learns the meaningful direction once data is clean.",
                Inches(0.5), Inches(6.6), Inches(5.5), Inches(0.85),
                fill=LIGHT, border=NAVY, size=11, bold=True)

    # Right: radial figure
    add_image(s, "results/figs/meeting/slide5_radial_distribution.png",
              Inches(6.5), Inches(1.6), width=Inches(6.5))
    add_textbox(s, "Radial distribution: benign (median 3.20) vs attacks (median 3.50). 0/13 inversions.",
                Inches(6.5), Inches(7.0), Inches(6.5), Inches(0.3),
                size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
#  SLIDE 6: HPS = C4 STATISTICALLY
# ============================================================
def slide_6(prs, fig_dir):
    s = blank_slide(prs)
    add_title(s, "6. Finding 2 — HPS Doesn't Beat Linear Probes", color=NAVY)
    add_subtitle(s,
                 "Geometric prior helps over flat (+0.049 TPR), but C4 still wins (+0.015 TPR)",
                 Inches(1.05))

    # Left: 4-method results table + cold-start
    headers = ["Method", "AUROC", "TPR@5%FPR", "Params"]
    rows = [
        ["MTP (Anthropic)",        "0.999", "0.995", "4,097"],
        ["C4 (linear probe)",      "0.999", "0.995", "4,097"],
        ["HPS (Lorentz)",          "0.997", "0.991", "262K"],
        ["HPS-Euclidean",          "0.968", "0.931", "262K"],
    ]
    add_table(s, headers, rows,
              Inches(0.5), Inches(1.6), Inches(5.5), Inches(2.0),
              body_size=12, bold_first_col=True, highlight_rows=[])

    # Key deltas
    add_textbox(s, "Key deltas:",
                Inches(0.5), Inches(3.8), Inches(5.5), Inches(0.3),
                size=13, color=NAVY, bold=True)
    add_bullets(s, [
        "HPS − HPS-Euclidean: +0.060 TPR (geometry helps over flat)",
        "HPS − C4: −0.004 TPR (linear probe still wins)",
        "MTP − C4: +0.000 (axis swap doesn't matter)",
    ], Inches(0.5), Inches(4.1), Inches(5.5), Inches(1.4), size=11)

    # Prediction agreement (the killer finding)
    add_callout(s,
                "Prediction agreement (calibrated 5% FPR thresholds):\n"
                "  • HPS catches that C4 misses:  0 examples\n"
                "  • C4 catches that HPS misses:  21 examples\n"
                "  • Pearson(HPS, C4) = 0.958\n"
                "  • OR-gate ensemble: no TPR gain, FPR doubles to 0.103\n\n"
                "→ HPS is essentially a noisy SUBSET of C4.",
                Inches(0.5), Inches(5.5), Inches(5.5), Inches(1.85),
                fill=RGBColor(0xFD, 0xED, 0xEC), border=RED, size=11, bold=True)

    # Right: 4-method bar chart on top, cold-start on bottom (DGX versions, cleaner)
    add_image(s, "results/figs/meeting/slide6_method_comparison.png",
              Inches(6.3), Inches(1.4), width=Inches(6.7))
    add_image(s, "results/figs/meeting/slide6_cold_start.png",
              Inches(6.3), Inches(4.7), width=Inches(6.7))


# ============================================================
#  SLIDE 7: VICUNA FAILURE
# ============================================================
def slide_7(prs, fig_dir):
    s = blank_slide(prs)
    add_title(s, "7. Finding 3 — Alignment-Mediated Failure on Vicuna", color=RED)
    add_subtitle(s,
                 "HPS catches 7.6% of GCG on Vicuna; C4 catches 99% — same code, only LLM differs",
                 Inches(1.05))

    # Left: the big number + per-attack table
    add_callout(s,
                "HPS GCG detection rate:\n"
                "  Llama-3-8B-Instruct (RLHF):  100% (172/172)\n"
                "  Vicuna-13B-v1.5 (SFT only):    7.6% (13/171)\n\n"
                "Cross-model gap: −92 percentage points",
                Inches(0.5), Inches(1.6), Inches(5.7), Inches(1.7),
                fill=RGBColor(0xFD, 0xED, 0xEC), border=RED, size=13, bold=True)

    # Per-attack on Vicuna
    add_textbox(s, "Per-attack on Vicuna-13B:",
                Inches(0.5), Inches(3.5), Inches(5.7), Inches(0.3),
                size=13, color=NAVY, bold=True)
    headers = ["Attack", "HPS", "C4", "Gap"]
    rows = [
        ["gcg",      "0.076", "0.994", "+0.92"],
        ["ijp",      "0.329", "0.933", "+0.60"],
        ["pair",     "0.389", "0.958", "+0.57"],
        ["puzzler",  "0.462", "1.000", "+0.54"],
        ["zulu",     "0.632", "1.000", "+0.37"],
        ["autodan",  "0.701", "1.000", "+0.30"],
        ["base64",   "0.922", "1.000", "+0.08"],
        ["drattack", "0.956", "1.000", "+0.04"],
        ["saa",      "0.994", "1.000", "+0.01"],
    ]
    add_table(s, headers, rows,
              Inches(0.5), Inches(3.85), Inches(5.7), Inches(2.65),
              body_size=10, bold_first_col=True, highlight_rows=[0, 1, 2, 3])

    add_textbox(s, "Mean detection: HPS = 0.61, C4 = 0.99",
                Inches(0.5), Inches(6.6), Inches(5.7), Inches(0.4),
                size=12, color=DARK, bold=True)

    # Right: figure (DGX version with built-in mechanism callout)
    add_image(s, "results/figs/meeting/slide7_vicuna_per_attack.png",
              Inches(6.5), Inches(1.4), width=Inches(6.5))

    add_callout(s,
                "Mechanism: RLHF concentrates harm features into compact regions; "
                "SFT-only leaves them diffuse. HPS's 12-feature compression preserves "
                "concentrated signals (Llama-3) but filters out diffuse ones (Vicuna).",
                Inches(6.5), Inches(5.7), Inches(6.5), Inches(1.6),
                fill=LIGHT, border=PURPLE, size=11, bold=False)


# ============================================================
#  SLIDE 8: BAILEY ET AL.
# ============================================================
def slide_8(prs, fig_dir):
    s = blank_slide(prs)
    add_title(s, "8. The Threat Surface — Bailey et al. (ICLR 2025)", color=NAVY)
    add_subtitle(s,
                 "All latent-space defenses fail under adaptive attacks — field-wide finding",
                 Inches(1.05))

    # Bailey et al. results table
    headers = ["Defense", "Standard recall", "Under adaptive attack", "Section"]
    rows = [
        ["Logistic regression probe",    "100%", "0%",                    "Sec 3.3"],
        ["MLP probe",                    "100%", "0%",                    "Sec 3.3"],
        ["Sparse autoencoder",           "91%",  "0%",                    "Sec 5"],
        ["OOD detector (Mahalanobis)",   "100%", "0%",                    "Sec 3.3"],
        ["Adv-trained probe (70 rounds)","perfect on past",
                                         "0% on new attacks",             "Sec 3.5"],
        ["Circuit Breakers",             "broken", "broken",              "Sec 3.4"],
        ["LAT (Latent Adv Training)",    "broken", "broken",              "Sec 3.4"],
    ]
    add_table(s, headers, rows,
              Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.6),
              body_size=12, bold_first_col=True,
              highlight_rows=[0, 1, 2, 3, 4])

    # Quote
    add_callout(s,
                '"Obfuscated activations are not rare exceptions but rather appear to be widespread '
                'in the latent space... an attacker can always find new activations that bypass the monitor."\n'
                "                                                                — Bailey et al. (ICLR 2025)",
                Inches(0.5), Inches(4.4), Inches(12.3), Inches(1.1),
                fill=LIGHT, border=NAVY, size=12)

    # Bottom line
    add_callout(s,
                "Bottom line for our paper:\n"
                "We do not claim adversarial robustness. Bailey et al. shows it doesn't exist for "
                "any activation-based probe variant, including ensembles and adversarially trained probes. "
                "We compare methods under standard (non-adaptive) conditions; "
                "adaptive robustness is a field-wide open problem.",
                Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.6),
                fill=RGBColor(0xFE, 0xF9, 0xE7), border=ORANGE, size=12, bold=True)


# ============================================================
#  SLIDE 9: OPEN QUESTIONS
# ============================================================
def slide_9(prs, fig_dir):
    s = blank_slide(prs)
    add_title(s, "9. What We Found vs. What's Open", color=NAVY)
    add_subtitle(s, "Need team input on framing, venue, and next experiments",
                 Inches(1.05))

    # Left: what we found
    add_textbox(s, "WHAT WE FOUND",
                Inches(0.5), Inches(1.6), Inches(6.3), Inches(0.4),
                size=15, color=GREEN, bold=True)
    add_textbox(s, "Methodology contributions",
                Inches(0.5), Inches(2.0), Inches(6.3), Inches(0.3),
                size=12, color=NAVY, bold=True)
    add_bullets(s, [
        "Length confound: AUROC=0.992 from length alone",
        "max_length confound: norm-only AUROC=1.000",
        "Train/test contamination: 15 prompts (1.15%) overlap",
    ], Inches(0.5), Inches(2.3), Inches(6.3), Inches(1.5), size=11)

    add_textbox(s, "Empirical findings",
                Inches(0.5), Inches(3.8), Inches(6.3), Inches(0.3),
                size=12, color=NAVY, bold=True)
    add_bullets(s, [
        "Geometric hypothesis CONFIRMED (0/13 inversions)",
        "HPS > HPS-Euclidean (+0.049 TPR matched-params)",
        "Vicuna alignment-mediated failure (gcg: 7.6% vs 99%)",
        "Linear probes (C4 ≈ MTP ≈ HPS) on aligned LLMs",
    ], Inches(0.5), Inches(4.1), Inches(6.3), Inches(2.0), size=11)

    # Right: open questions
    add_textbox(s, "WHAT'S OPEN — DECISIONS FOR THE TEAM",
                Inches(7.0), Inches(1.6), Inches(6.0), Inches(0.4),
                size=15, color=ORANGE, bold=True)
    add_textbox(s, ["Q1: Where to publish?",
                    "    TMLR (60-65%) | AAAI (35-50%) | USENIX (25-40%)",
                    "    Recommendation: TMLR",
                    "",
                    "Q2: Run adaptive attacks experiment?",
                    "    1-2 days compute, predicted: confirms Bailey et al.",
                    "",
                    "Q3: Add Llama-2-7b-chat? CRITICAL.",
                    "    SFT+RLHF on Vicuna's base model = clean ablation",
                    "    of RLHF as the variable (vs Vicuna SFT-only)",
                    "",
                    "Q4: What's the lead?",
                    "    methodology critique | architectural | alignment-mediated",
                    "",
                    "Q5: Multi-turn pivot? See backup slide for details.",
                    "    6mo vs 6wk TMLR submission"],
                Inches(7.0), Inches(2.0), Inches(6.0), Inches(5.0),
                size=11, color=DARK)

    # Bottom callout
    add_callout(s,
                "Looking for input on: paper venue, framing priority, "
                "whether to add adaptive attacks, and whether to pivot to multi-turn.",
                Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6),
                fill=RGBColor(0xFE, 0xF9, 0xE7), border=ORANGE, size=12, bold=True)


# ============================================================
#  SLIDE 10: BACKUP — Q&A REFERENCE
# ============================================================
def slide_qa_backup(prs, fig_dir):
    s = blank_slide(prs)
    add_title(s, "Backup: Anticipated Q&A", color=GRAY, size=28)

    qas = [
        ("Q: If geometric priors don't beat linear probes, why publish?",
         "A: Three field-wide methodology issues + controlled comparison + alignment-mediated failure."),
        ("Q: Why no adversarial robustness?",
         "A: Bailey et al. (ICLR 2025) shows all latent probes fail. Confirmed limitation."),
        ("Q: Are JBShield/HSF/RTV wrong?",
         "A: No — but their benchmarks have inherited methodology issues. Fix is diverse benign."),
        ("Q: Why Vicuna-13B?",
         "A: Same family as Llama-2 base, but SFT-only. Isolates alignment as the variable."),
        ("Q: What if HPS fails on Llama-2-chat too?",
         "A: Sharpens story to 'specific aligned-model class'. Either way, ablation strengthens paper."),
        ("Q: How does this compare to JBShield's 99% detection?",
         "A: Their 99% likely includes substantial length signal. Reproducing on diverse benign would tell us."),
        ("Q: TMLR timeline?",
         "A: 4-6 weeks writing → submission August → decision November (3 month review)."),
        ("Q: Compute budget used?",
         "A: ~10h cache extractions + 5h experiments. ~15 GPU-days project total."),
    ]

    y = Inches(1.5)
    for q, a in qas:
        add_textbox(s, q, Inches(0.5), y, Inches(12.3), Inches(0.3),
                    size=12, color=NAVY, bold=True)
        add_textbox(s, a, Inches(0.7), y + Inches(0.3), Inches(12.0), Inches(0.4),
                    size=11, color=DARK)
        y += Inches(0.7)


# ============================================================
#  Build
# ============================================================
def main():
    p = argparse.ArgumentParser()
    p.add_argument("--output", default="team_meeting.pptx")
    p.add_argument("--figures", default="figures_for_meeting",
                   help="Directory with figure PNGs")
    args = p.parse_args()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print(f"\n  Building slides → {args.output}")
    print(f"  Figure directory: {args.figures}")

    slide_title(prs)
    slide_1(prs, args.figures)
    slide_2(prs, args.figures)
    slide_3(prs, args.figures)
    slide_4(prs, args.figures)
    slide_5(prs, args.figures)
    slide_6(prs, args.figures)
    slide_7(prs, args.figures)
    slide_8(prs, args.figures)
    slide_9(prs, args.figures)
    slide_qa_backup(prs, args.figures)

    prs.save(args.output)
    print(f"  Saved {len(prs.slides)} slides → {args.output}")
    size_kb = os.path.getsize(args.output) // 1024
    print(f"  File size: {size_kb} KB\n")


if __name__ == "__main__":
    main()
